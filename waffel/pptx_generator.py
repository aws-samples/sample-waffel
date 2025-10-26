#!/usr/bin/env python3
"""PowerPoint presentation generator for WAFFEL reports."""

import io
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .config import WAFFELConfig

class PowerPointGenerator:
    """Generates PowerPoint presentations from Well-Architected data"""

    def __init__(self):
        self.config = WAFFELConfig()

    def generate(self, pillars, workload_props, output_path):
        """Generate PowerPoint file with improvement items as colored rectangles"""

        prs = Presentation()

        # Create title slide
        self._create_title_slide(prs, workload_props)

        # Create improvement items slides
        self._create_improvement_slides(prs, pillars)

        # Add Eisenhower matrix slide
        self._create_eisenhower_slide(prs)

        # Add Priority slide
        self._create_priority_slide(prs)

        prs.save(output_path)

    def _create_title_slide(self, prs, workload_props):
        """Create title slide with workload information"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        workload_name = workload_props.get('Workload name', 'Well-Architected Assessment')
        title.text = f"{workload_name}"
        subtitle.text = "Improvement Plan Overview"

    def _create_improvement_slides(self, prs, pillars):
        """Create single slide with all improvement items as small colored rectangles"""

        # Collect all improvement items
        improvement_items = []
        for pillar_name, questions in pillars.items():
            for question in questions:
                if question.get('improvement_items'):
                    for item_data in question['improvement_items']:
                        improvement_items.append({
                            'pillar': pillar_name,
                            'item': item_data['item'],
                            'url': item_data['url'],
                            'question_id': question['question_id'],
                            'risk_level': question.get('risk_level', '')
                        })

        if not improvement_items:
            # Create slide with "No improvement items" message
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Improvement Items"

            content = slide.placeholders[1]
            content.text = "No improvement items found in this assessment."
            return

        # Create single slide with all items
        self._create_compact_items_slide(prs, improvement_items)

    def _create_compact_items_slide(self, prs, items):
        """Create a single slide with small improvement item rectangles"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Improvement Items Overview"
        title_frame.paragraphs[0].font.size = Inches(0.3)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Small rectangle dimensions (3x1 cm)
        rect_width = Inches(1.18)   # 3 cm
        rect_height = Inches(0.39)  # 1 cm
        margin = Inches(0.1)

        # Calculate grid layout to fit all items
        start_x = Inches(0.5)
        start_y = Inches(1.2)

        # Calculate how many columns we can fit
        available_width = Inches(9)
        cols = int((available_width - start_x) / (rect_width + margin))

        for idx, item in enumerate(items):
            row = idx // cols
            col = idx % cols

            x = start_x + col * (rect_width + margin)
            y = start_y + row * (rect_height + margin)

            # Create rectangle
            rect = slide.shapes.add_shape(
                1,  # Rectangle shape type
                x, y, rect_width, rect_height
            )

            # Set pillar color (pastel)
            pillar_color = self.config.PILLAR_COLORS.get(item['pillar'], (200, 200, 200))
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(*pillar_color)

            # Add border - same color for all (no special high risk border)
            rect.line.color.rgb = RGBColor(100, 100, 100)  # Gray border for all
            rect.line.width = Inches(0.01)  # Thin border

            # Add text with question ID prefix
            text_frame = rect.text_frame
            text_frame.margin_left = Inches(0.02)
            text_frame.margin_right = Inches(0.02)
            text_frame.margin_top = Inches(0.02)
            text_frame.margin_bottom = Inches(0.02)
            text_frame.word_wrap = True

            # Question ID (bold) + improvement text in same string
            p = text_frame.paragraphs[0]
            item_text = item['item']
            # Calculate space for question ID + space + text
            question_id = item['question_id']
            available_chars = 50 - len(question_id) - 1  # Reserve space for ID + space
            if len(item_text) > available_chars:
                item_text = item_text[:available_chars-3] + "..."

            # Create text with question ID and improvement text
            p = text_frame.paragraphs[0]
            p.clear()

            # Add question ID as bold hyperlink
            question_run = p.add_run()
            question_run.text = question_id + " "
            question_run.font.bold = True
            question_run.font.size = Pt(8)  # 8pt font

            # Add hyperlink to question ID if URL exists
            if item['url']:
                question_run.hyperlink.address = item['url']

            # Force text color to black (must be after hyperlink creation)
            question_run.font.color.rgb = RGBColor(0, 0, 0)

            # Add improvement text as normal
            text_run = p.add_run()
            text_run.text = item_text
            text_run.font.size = Pt(8)  # 8pt font
            text_run.font.color.rgb = RGBColor(0, 0, 0)

    def _create_eisenhower_slide(self, prs):
        """Create slide with full-height Eisenhower matrix image"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Prioritization Framework"
        title_frame.paragraphs[0].font.size = Inches(0.25)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Download and add image (full height)
        try:
            image_url = "https://docs.aws.amazon.com/images/wellarchitected/latest/userguide/images/eisenhower.png"
            response = requests.get(image_url, timeout=10)  # nosec B113
            response.raise_for_status()

            # Add image to slide - full height (6.5 inches available after title)
            image_stream = io.BytesIO(response.content)
            slide.shapes.add_picture(image_stream, Inches(0.5), Inches(0.9), height=Inches(6.5))

        except Exception:
            # If image download fails, add text instead
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            text_frame = text_box.text_frame
            text_frame.text = "Eisenhower Matrix for prioritizing improvement items:\n\n• Urgent + Important = Do First\n• Important + Not Urgent = Schedule\n• Urgent + Not Important = Delegate\n• Not Urgent + Not Important = Eliminate"
            text_frame.paragraphs[0].font.size = Pt(16)

    def _create_priority_slide(self, prs):
        """Create Priority slide with empty list"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = "Priority"

        # Add empty bullet list
        content = slide.placeholders[1]
        content.text = "\n\n\n\n"
